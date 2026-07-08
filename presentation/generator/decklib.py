"""decklib — a small, opinionated framework over python-pptx for building a
WSO2-branded, 16:9 enterprise slide deck with consistent geometry.

Design goals: correct-by-construction geometry (fixed grid), flat modern styling
(no theme shadows / stray table styles), and reusable diagram primitives so the
same node/arrow look is used everywhere.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Palette (WSO2-branded, light enterprise theme)
# ----------------------------------------------------------------------------
def C(h):
    return RGBColor.from_string(h)

ORANGE    = "FF7300"   # WSO2 primary
ORANGE_DK = "C85A00"
ORANGE_LT = "FFEBDA"
INK       = "1B2733"   # headings / body
INK_SOFT  = "3E4C59"   # secondary text
GRAY      = "6B7885"
GRAY_MD   = "9AA5B1"
GRAY_LT   = "F4F6F8"
PANEL     = "F6F8FA"
PANEL_2   = "EEF2F6"
LINE      = "DCE3E9"
WHITE     = "FFFFFF"
NAVY      = "15222E"   # dark backgrounds (title / section)
NAVY_2    = "20374A"   # faint accents on navy
NAVY_TXT  = "AEBECD"   # muted text on navy

TEAL      = "0E8FA8"   # Solution 1 / Ballerina accent
TEAL_DK   = "0A6E82"
TEAL_LT   = "E0F1F5"
VIOLET    = "6A4CE0"   # Solution 2 / LangChain accent
VIOLET_DK = "5138C0"
VIOLET_LT = "ECE8FC"
SPLUNK    = "5AA82B"   # splunk green
SPLUNK_LT = "EAF4E1"
DATADOG   = "7A3FE4"   # datadog purple
DATADOG_LT= "EFE8FC"
GREEN     = "2E9E5B"
GREEN_LT  = "E4F4EB"
AMBER     = "D98B22"
AMBER_LT  = "FBEFD9"
RED       = "CE4A45"
RED_LT    = "FBE7E6"

FONT = "Arial"
MONO = "Consolas"

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN  = 0.62
CW      = SLIDE_W - 2 * MARGIN     # content width
BODY_TOP = 1.9
BODY_BOT = 7.02
BODY_H   = BODY_BOT - BODY_TOP

DECK_TAG = "DevOps OverSight Agent  ·  POC Architecture Review"

# ----------------------------------------------------------------------------
# Deck / slide scaffolding
# ----------------------------------------------------------------------------
def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def add_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    if bg:
        set_bg(slide, bg)
    return slide


def set_bg(slide, h):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
    r.fill.solid(); r.fill.fore_color.rgb = C(h)
    r.line.fill.background()
    r.shadow.inherit = False
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return r


# ----------------------------------------------------------------------------
# Primitive shapes
# ----------------------------------------------------------------------------
def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line); sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf


def _apply_runs(p, runs, size, color, bold, font):
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, ov in runs:
        r = p.add_run(); r.text = text
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        it = ov.get("italic", None)
        if it is not None:
            f.italic = it
        f.color.rgb = C(ov.get("color", color))


def _apply_bullet(p, char, color_hex, level=0):
    pPr = p._p.get_or_add_pPr()
    marL = int((0.28 + 0.30 * level) * 914400)
    pPr.set("marL", str(marL))
    pPr.set("indent", str(-int(0.28 * 914400)))
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": color_hex})
    buClr.append(srgb)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buClr); pPr.append(buFont); pPr.append(buChar)


def para(tf, runs, first=False, size=14, color=INK, bold=False, font=FONT,
         align=PP_ALIGN.LEFT, before=0, after=6, spacing=1.0, level=0,
         bullet=None, bullet_color=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if spacing is not None:
        p.line_spacing = spacing
    if before is not None:
        p.space_before = Pt(before)
    if after is not None:
        p.space_after = Pt(after)
    if bullet:
        _apply_bullet(p, bullet, bullet_color or ORANGE, level)
    _apply_runs(p, runs, size, color, bold, font)
    return p


# ----------------------------------------------------------------------------
# Brand marks & chrome
# ----------------------------------------------------------------------------
def wso2_mark(slide, x, y, color=ORANGE, size=18):
    tb, tf = textbox(slide, x, y, 2.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("wso", {"bold": True, "color": color, "size": size}),
              ("2", {"bold": True, "color": ORANGE, "size": size})], first=True)
    return tb


def footer(slide, page_no, total=None, tag=DECK_TAG):
    tb, tf = textbox(slide, MARGIN, 7.12, CW - 1.4, 0.28, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, tag, first=True, size=8, color=GRAY_MD)
    tb2, tf2 = textbox(slide, SLIDE_W - MARGIN - 1.4, 7.12, 1.4, 0.28, anchor=MSO_ANCHOR.MIDDLE)
    txt = f"{page_no:02d}" + (f"  /  {total:02d}" if total else "")
    para(tf2, txt, first=True, size=8, color=GRAY_MD, align=PP_ALIGN.RIGHT)


def kicker_title(slide, kicker, title, accent=ORANGE, title_size=27):
    """Standard content-slide header: eyebrow + title + rule."""
    rect(slide, MARGIN, 0.56, 0.34, 0.10, fill=accent)
    tb, tf = textbox(slide, MARGIN + 0.46, 0.44, CW - 0.46, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, kicker.upper(), first=True, size=11, color=accent, bold=True)
    tb2, tf2 = textbox(slide, MARGIN, 0.82, CW, 0.86, anchor=MSO_ANCHOR.TOP)
    para(tf2, title, first=True, size=title_size, color=INK, bold=True, spacing=1.0)
    rect(slide, MARGIN, 1.66, CW, 0.014, fill=LINE)


def content_slide(prs, kicker, title, accent=ORANGE, page=None, total=None):
    slide = add_slide(prs, WHITE)
    kicker_title(slide, kicker, title, accent)
    if page is not None:
        footer(slide, page, total)
    return slide


def section_slide(prs, number, title, subtitle=None, accent=ORANGE):
    slide = add_slide(prs, NAVY)
    # giant faint section number, top-right
    tb, tf = textbox(slide, SLIDE_W - 5.4, -0.35, 5.0, 4.2, anchor=MSO_ANCHOR.TOP)
    para(tf, number, first=True, size=210, color=NAVY_2, bold=True, align=PP_ALIGN.RIGHT)
    rect(slide, MARGIN, 3.02, 0.9, 0.11, fill=accent)
    tb2, tf2 = textbox(slide, MARGIN, 3.28, 10.5, 0.4)
    para(tf2, "SECTION", first=True, size=12.5, color=accent, bold=True)
    tb3, tf3 = textbox(slide, MARGIN, 3.66, 11.2, 1.7)
    para(tf3, title, first=True, size=40, color=WHITE, bold=True, spacing=1.02)
    if subtitle:
        tb4, tf4 = textbox(slide, MARGIN, 5.4, 10.6, 1.1)
        para(tf4, subtitle, first=True, size=15, color=NAVY_TXT, spacing=1.15)
    wso2_mark(slide, MARGIN, 6.86, color=WHITE, size=16)
    return slide


# ----------------------------------------------------------------------------
# Cards, panels, callouts, bullets
# ----------------------------------------------------------------------------
def panel(slide, x, y, w, h, fill=PANEL, line=LINE, line_w=1.0, radius=0.045):
    return rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)


def card_header(slide, x, y, w, title, accent=ORANGE, tsize=13.5, icon=None):
    """A left accent tick + bold title used at the top of a panel."""
    rect(slide, x, y + 0.02, 0.09, 0.26, fill=accent)
    tb, tf = textbox(slide, x + 0.22, y - 0.02, w - 0.22, 0.34, anchor=MSO_ANCHOR.MIDDLE)
    lead = f"{icon}  " if icon else ""
    para(tf, (lead + title), first=True, size=tsize, color=INK, bold=True)


def bullets(slide, x, y, w, h, items, size=13, gap=7, color=INK,
            bullet_color=ORANGE, anchor=MSO_ANCHOR.TOP, lead_color=INK,
            spacing=1.06):
    """items: list of dicts. Each: {lead, text, level, char, bc, color}.
    Or a plain string. Lead is rendered bold in lead_color."""
    tb, tf = textbox(slide, x, y, w, h, anchor=anchor)
    for i, it in enumerate(items):
        if isinstance(it, str):
            it = {"text": it}
        lvl = it.get("level", 0)
        char = it.get("char", "▪" if lvl == 0 else "–")
        bc = it.get("bc", bullet_color if lvl == 0 else GRAY_MD)
        runs = []
        if it.get("lead"):
            runs.append((it["lead"] + "  ", {"bold": True, "color": it.get("lead_color", lead_color)}))
        if it.get("text"):
            runs.append((it["text"], {"color": it.get("color", color)}))
        if not runs:
            runs = [(" ", {})]
        para(tf, runs, first=(i == 0), size=it.get("size", size), color=color,
             align=PP_ALIGN.LEFT, before=0, after=it.get("after", gap),
             spacing=spacing, level=lvl, bullet=char, bullet_color=bc)
    return tf


def stat(slide, x, y, w, big, label, accent=ORANGE, big_size=40, lab_size=11,
         sub=None, h=1.5, fill=PANEL, line=LINE):
    panel(slide, x, y, w, h, fill=fill, line=line)
    tb, tf = textbox(slide, x + 0.12, y + 0.12, w - 0.24, h - 0.24, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, first=True, size=big_size, color=accent, bold=True, after=2, align=PP_ALIGN.CENTER)
    para(tf, label, size=lab_size, color=INK, bold=True, align=PP_ALIGN.CENTER, after=0, spacing=1.02)
    if sub:
        para(tf, sub, size=9, color=GRAY, align=PP_ALIGN.CENTER, before=1, after=0, spacing=1.02)


def callout(slide, x, y, w, h, title, text, accent=ORANGE, fill=ORANGE_LT,
            tsize=12, bsize=10.5, icon="◆"):
    panel(slide, x, y, w, h, fill=fill, line=None)
    rect(slide, x, y, 0.10, h, fill=accent)  # left rail
    tb, tf = textbox(slide, x + 0.26, y + 0.14, w - 0.42, h - 0.28, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, f"{icon}  {title}", first=True, size=tsize, color=accent, bold=True, after=3, spacing=1.0)
    para(tf, text, size=bsize, color=INK_SOFT, spacing=1.12, after=0)


# ----------------------------------------------------------------------------
# Diagram primitives
# ----------------------------------------------------------------------------
def node(slide, x, y, w, h, title, sub=None, fill=WHITE, line=LINE, line_w=1.4,
         tcolor=INK, scolor=GRAY, tsize=12, ssize=9, radius=0.10, tag=None,
         tag_color=None, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    box = rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)
    tb, tf = textbox(slide, x + 0.08, y + 0.04, w - 0.16, h - 0.08, anchor=anchor)
    first = True
    if tag:
        para(tf, tag.upper(), first=True, size=7.5, color=tag_color or GRAY, bold=True,
             align=align, after=2, spacing=1.0)
        first = False
    para(tf, title, first=first, size=tsize, color=tcolor, bold=True, align=align,
         after=(2 if sub else 0), spacing=1.0)
    if sub:
        para(tf, sub, size=ssize, color=scolor, align=align, after=0, spacing=1.02)
    return box


def connect(slide, x1, y1, x2, y2, color=GRAY, width=1.6, dashed=False,
            head=True, tail=False):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = C(color); cxn.line.width = Pt(width)
    ln = cxn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if tail:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if head:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    cxn.shadow.inherit = False
    return cxn


def elbow(slide, x1, y1, x2, y2, color=GRAY, width=1.6, dashed=False,
          head=True, tail=False):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW,
                                     Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = C(color); cxn.line.width = Pt(width)
    ln = cxn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if tail:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if head:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    cxn.shadow.inherit = False
    return cxn


def edge_label(slide, cx, cy, text, w=1.7, size=8.5, color=INK_SOFT, bold=False,
               fill=WHITE, align=PP_ALIGN.CENTER, h=0.26):
    if fill:
        rect(slide, cx - w / 2, cy - h / 2, w, h, fill=fill, line=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.4)
    tb, tf = textbox(slide, cx - w / 2, cy - h / 2, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, first=True, size=size, color=color, bold=bold, align=align)


def chip(slide, x, y, text, w=None, color=ORANGE, fill=ORANGE_LT, size=9.5,
         bold=True, h=0.30):
    if w is None:
        w = 0.20 + 0.085 * len(text)
    rect(slide, x, y, w, h, fill=fill, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tb, tf = textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, first=True, size=size, color=color, bold=bold, align=PP_ALIGN.CENTER)
    return w


def numbered_step(slide, x, y, w, n, title, text, accent=ORANGE, h=0.92,
                  tsize=11.5, bsize=9.5):
    """A numbered step chip + title + description inside a light panel."""
    panel(slide, x, y, w, h, fill=WHITE, line=LINE)
    d = 0.42
    rect(slide, x + 0.16, y + 0.16, d, d, fill=accent, shape=MSO_SHAPE.OVAL)
    tb, tf = textbox(slide, x + 0.16, y + 0.16, d, d, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, str(n), first=True, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb2, tf2 = textbox(slide, x + 0.74, y + 0.10, w - 0.86, h - 0.2, anchor=MSO_ANCHOR.MIDDLE)
    para(tf2, title, first=True, size=tsize, color=INK, bold=True, after=2, spacing=1.0)
    para(tf2, text, size=bsize, color=INK_SOFT, spacing=1.08, after=0)


# ----------------------------------------------------------------------------
# Tables
# ----------------------------------------------------------------------------
_NO_STYLE_NO_GRID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def make_table(slide, x, y, w, col_fracs, data, header=True, fsize=10,
               hsize=10.5, header_fill=INK, header_color=WHITE, row_h=0.42,
               header_h=0.44, zebra=PANEL, first_col_bold=False, align0=PP_ALIGN.LEFT):
    nrows = len(data); ncols = len(col_fracs)
    total_h = header_h + row_h * (nrows - 1)
    gt = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                                Inches(w), Inches(total_h)).table
    # neutralize built-in table style
    tblPr = gt._tbl.tblPr
    for el in list(tblPr):
        if el.tag == qn("a:tableStyleId"):
            tblPr.remove(el)
    sid = tblPr.makeelement(qn("a:tableStyleId"), {})
    sid.text = _NO_STYLE_NO_GRID
    tblPr.append(sid)
    tblPr.set("firstRow", "0"); tblPr.set("bandRow", "0")

    for i, f in enumerate(col_fracs):
        gt.columns[i].width = Inches(w * f)
    for r in range(nrows):
        gt.rows[r].height = Inches(header_h if (r == 0 and header) else row_h)
        for c in range(ncols):
            cell = gt.cell(r, c)
            spec = data[r][c]
            if not isinstance(spec, dict):
                spec = {"text": str(spec)}
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.10); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            if r == 0 and header:
                cell.fill.solid(); cell.fill.fore_color.rgb = C(spec.get("fill", header_fill))
            else:
                fill = spec.get("fill")
                if fill:
                    cell.fill.solid(); cell.fill.fore_color.rgb = C(fill)
                elif zebra and (r % 2 == 1):
                    cell.fill.solid(); cell.fill.fore_color.rgb = C(zebra)
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = C(WHITE)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            default_align = align0 if c == 0 else PP_ALIGN.LEFT
            p.alignment = spec.get("align", default_align)
            p.line_spacing = 1.0
            rr = p.add_run(); rr.text = spec.get("text", "")
            fnt = rr.font
            fnt.name = FONT
            fnt.size = Pt(spec.get("size", hsize if (r == 0 and header) else fsize))
            is_hdr = (r == 0 and header)
            fnt.bold = spec.get("bold", is_hdr or (c == 0 and first_col_bold))
            fnt.color.rgb = C(spec.get("color", header_color if is_hdr else INK))
    return gt


# ----------------------------------------------------------------------------
# QA: report any shape that spills off the slide
# ----------------------------------------------------------------------------
def qa_report(prs):
    issues = []
    W = Inches(SLIDE_W); H = Inches(SLIDE_H); tol = Inches(0.06)
    for idx, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
                if l is None or t is None or w is None or h is None:
                    continue
            except Exception:
                continue
            if l < -tol or t < -tol or (l + w) > W + tol or (t + h) > H + tol:
                nm = getattr(sh, "name", "?")
                txt = ""
                if sh.has_text_frame:
                    txt = sh.text_frame.text[:40].replace("\n", " ")
                issues.append((idx, nm, round(l / 914400, 2), round(t / 914400, 2),
                               round((l + w) / 914400, 2), round((t + h) / 914400, 2), txt))
    return issues
